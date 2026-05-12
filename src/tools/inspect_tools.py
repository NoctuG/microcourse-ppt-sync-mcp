"""Inspection tools for MCP."""

import os
import json
from pathlib import Path
from typing import Dict, Any
from src.utils import get_logger, validators


logger = get_logger(__name__)


def inspect_project(project_path: str) -> Dict[str, Any]:
    """Inspect project directory structure.
    
    Args:
        project_path: Path to project directory
        
    Returns:
        Dictionary with project information
    """
    try:
        project_path = str(validators.validate_project_path(project_path))
        
        result = {
            "status": "success",
            "project_path": project_path,
            "structure": {},
            "files": {},
        }
        
        # Check standard directories
        standard_dirs = ["input", "work", "output"]
        for dir_name in standard_dirs:
            dir_path = Path(project_path) / dir_name
            result["structure"][dir_name] = {
                "exists": dir_path.exists(),
                "is_dir": dir_path.is_dir(),
            }
            
            if dir_path.exists():
                files = list(dir_path.glob("*"))
                result["files"][dir_name] = [f.name for f in files]
        
        # Check for key files
        key_files = {
            "input": ["lesson.pptx"],
            "work": ["timeline.json", "transcript.json"],
            "output": ["lesson_timed.pptx", "ppt_bg.mp4", "sync_report.md"],
        }
        
        result["key_files"] = {}
        for dir_name, files in key_files.items():
            result["key_files"][dir_name] = {}
            for file_name in files:
                file_path = Path(project_path) / dir_name / file_name
                result["key_files"][dir_name][file_name] = file_path.exists()
        
        logger.info(f"Inspected project: {project_path}")
        return result
    except Exception as e:
        logger.error(f"Failed to inspect project: {e}")
        return {
            "status": "error",
            "message": str(e),
        }


def inspect_ppt(ppt_path: str) -> Dict[str, Any]:
    """Inspect PPT file.
    
    Args:
        ppt_path: Path to PPT file
        
    Returns:
        Dictionary with PPT information
    """
    try:
        ppt_path = str(validators.validate_ppt_file(ppt_path))
        
        result = {
            "status": "success",
            "ppt_path": ppt_path,
            "file_size": os.path.getsize(ppt_path),
            "file_format": Path(ppt_path).suffix,
        }
        
        # Try to get slide count using python-pptx if available
        try:
            from pptx import Presentation
            prs = Presentation(ppt_path)
            result["slide_count"] = len(prs.slides)
            
            # Get basic info about each slide
            result["slides"] = []
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    "index": i,
                    "shape_count": len(slide.shapes),
                    "has_notes": bool(slide.notes.text),
                }
                result["slides"].append(slide_info)
        except ImportError:
            logger.warning("python-pptx not available - limited PPT inspection")
            result["slide_count"] = "unknown"
            result["note"] = "Install python-pptx for detailed inspection"
        
        logger.info(f"Inspected PPT: {ppt_path}")
        return result
    except Exception as e:
        logger.error(f"Failed to inspect PPT: {e}")
        return {
            "status": "error",
            "message": str(e),
        }
